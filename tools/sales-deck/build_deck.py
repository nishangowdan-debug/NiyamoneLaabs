"""
Niyamone Labs × Sree Diagnostics — client walkthrough deck generator.

Run:
    pip install python-pptx
    python tools/sales-deck/build_deck.py

Output:
    tools/sales-deck/Niyamone-SreeDiagnostics-Walkthrough.pptx

Drop screenshots into tools/sales-deck/screenshots/ using the filenames
listed in SCREENSHOTS.md — the script will pick them up automatically.
If a screenshot is missing it draws a dashed placeholder with the
expected filename so you see exactly what is still to be captured.
"""

from __future__ import annotations
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand tokens (mirrors src/styles.css) ─────────────────────────────
BRAND       = RGBColor(0x0E, 0x4F, 0x8C)
BRAND_DEEP  = RGBColor(0x0A, 0x3A, 0x6B)
ACCENT      = RGBColor(0x00, 0xC3, 0xFF)
GOLD        = RGBColor(0xC9, 0xA2, 0x4B)
INK         = RGBColor(0x0F, 0x1B, 0x2D)
INK_SOFT    = RGBColor(0x2A, 0x37, 0x4A)
INK_MUTED   = RGBColor(0x65, 0x75, 0x8C)
INK_FAINT   = RGBColor(0x99, 0xA6, 0xB8)
SURFACE     = RGBColor(0xF4, 0xF7, 0xFB)
SURFACE_SUB = RGBColor(0xED, 0xF1, 0xF7)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GOOD_FG     = RGBColor(0x11, 0x7A, 0x3A)
WARN_FG     = RGBColor(0x8B, 0x5A, 0x0F)

ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
OUT = ROOT / "Niyamone-SreeDiagnostics-Walkthrough.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────
def set_solid(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill: RGBColor | None = None, line: RGBColor | None = None, line_w: float = 0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=15, color=INK_SOFT, gap=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "▸  " + b
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_pill(slide, x, y, label, *, fill=BRAND, fg=WHITE, w=Inches(1.8), h=Inches(0.32)):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Calibri"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = fg
    return pill


def brand_top_bar(slide):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=BRAND)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(0.04), fill=ACCENT)


def page_footer(slide, idx, total):
    add_text(slide, Inches(0.6), Inches(7.1), Inches(8.0), Inches(0.3),
             "Niyamone Labs  ×  Sree Diagnostics  ·  Confidential client walkthrough",
             size=9, color=INK_FAINT)
    add_text(slide, Inches(12.0), Inches(7.1), Inches(0.9), Inches(0.3),
             f"{idx} / {total}", size=9, color=INK_FAINT, align=PP_ALIGN.RIGHT)


def add_screenshot(slide, name: str, x, y, w, h, caption: str | None = None):
    """Embed screenshots/<name> if present; otherwise draw a labelled placeholder."""
    path = SHOTS / name
    border = add_rect(slide, x, y, w, h, fill=WHITE, line=BRAND, line_w=1.5)
    if path.exists():
        slide.shapes.add_picture(str(path), x, y, w, h)
    else:
        # Dashed inner placeholder
        inner = add_rect(slide, x + Inches(0.08), y + Inches(0.08),
                         w - Inches(0.16), h - Inches(0.16),
                         fill=SURFACE_SUB, line=INK_FAINT, line_w=0.75)
        ln = inner.line.color  # noqa: F841
        # mark line dashed
        sppr = inner.line._get_or_add_ln()
        prst = etree.SubElement(sppr, qn("a:prstDash"))
        prst.set("val", "dash")
        add_text(slide, x, y + h/2 - Inches(0.35), w, Inches(0.35),
                 f"⤓  Drop screenshot:  {name}",
                 size=13, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + h/2, w, Inches(0.3),
                 "tools/sales-deck/screenshots/" + name,
                 size=10, color=INK_MUTED, align=PP_ALIGN.CENTER)
    if caption:
        add_text(slide, x, y + h + Inches(0.05), w, Inches(0.35),
                 caption, size=10, color=INK_MUTED,
                 align=PP_ALIGN.CENTER, font="Calibri")


# ── Slide builders ─────────────────────────────────────────────────────
def make_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=BRAND_DEEP)
    add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(0.12), fill=ACCENT)
    add_text(s, Inches(0.8), Inches(0.7), Inches(6), Inches(0.4),
             "NIYAMONE LABS", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), Inches(2.3), Inches(12), Inches(1.4),
             "The diagnostic-lab\noperating system,\nbuilt for India.",
             size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(5.0), Inches(12), Inches(0.5),
             "Walkthrough for  ·  Sree Diagnostics, Vijayawada",
             size=20, color=WHITE)
    add_text(s, Inches(0.8), Inches(6.4), Inches(12), Inches(0.4),
             "Confidential  ·  Prepared for client visit  ·  May 2026",
             size=11, color=INK_FAINT)
    return s


def make_value_prop(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SURFACE)
    brand_top_bar(s)
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             "WHY WE'RE IN THIS ROOM", size=11, bold=True, color=BRAND)
    add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(1.6),
             "Run every branch of Sree Diagnostics\nfrom one screen — and prove it on day one.",
             size=36, bold=True, color=INK)
    chips = [
        ("FASTER",   "Report TAT cut by 40 %",         GOOD_FG),
        ("CLEANER",  "Zero paperwork at billing",      BRAND),
        ("SAFER",    "Every discount audit-trailed",   WARN_FG),
    ]
    base_x = Inches(0.6)
    for i, (k, v, c) in enumerate(chips):
        gx = base_x + Inches(i * 4.1)
        add_rect(s, gx, Inches(4.0), Inches(3.9), Inches(2.4),
                 fill=WHITE, line=INK_FAINT, line_w=0.5)
        add_pill(s, gx + Inches(0.3), Inches(4.2), k, fill=c, w=Inches(1.1))
        add_text(s, gx + Inches(0.3), Inches(4.9), Inches(3.5), Inches(1.6),
                 v, size=22, bold=True, color=INK)
    page_footer(s, idx, total)
    return s


def make_problem(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SURFACE)
    brand_top_bar(s)
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             "THE PROBLEM", size=11, bold=True, color=BRAND)
    add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.8),
             "What every 5-branch lab in India tells us.",
             size=32, bold=True, color=INK)
    items = [
        ("63 %",  "of report errors trace back to manual transcription between paper, Excel, and the LIMS."),
        ("3.4 hrs", "lost per receptionist per day re-entering the same patient into 4 different systems."),
        ("₹ 4.2 L", "average annual leakage from un-tracked cash-counter discounts and refunds."),
    ]
    y = Inches(2.8)
    for stat, copy in items:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.2),
                 fill=WHITE, line=INK_FAINT, line_w=0.5)
        add_text(s, Inches(0.9), y + Inches(0.2), Inches(2.4), Inches(0.9),
                 stat, size=32, bold=True, color=BRAND)
        add_text(s, Inches(3.6), y + Inches(0.3), Inches(9.0), Inches(0.8),
                 copy, size=16, color=INK_SOFT)
        y += Inches(1.35)
    page_footer(s, idx, total)
    return s


def make_promise(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SURFACE)
    brand_top_bar(s)
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             "THE NIYAMONE PROMISE", size=11, bold=True, color=BRAND)
    add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(1.4),
             "One login. Every branch.\nEvery workflow. Live data.",
             size=36, bold=True, color=INK)
    add_bullets(s, Inches(0.6), Inches(3.6), Inches(5.8), Inches(3),
                ["Reception, lab, billing, pharmacy, IPD — same app, same login.",
                 "Branch switch in one click. No re-login. No data re-sync.",
                 "Your data lives in *your* Supabase tenant — full export, any time.",
                 "Indian-first: GST, TDS 194J, UHID, IFSC, PAN — built in, not bolted on."],
                size=16)
    add_screenshot(s, "00-dashboard-hero.png",
                   Inches(6.8), Inches(3.4), Inches(6.0), Inches(3.4),
                   "Operational dashboard — one branch or all five, at a glance.")
    page_footer(s, idx, total)
    return s


def make_outcomes(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SURFACE)
    brand_top_bar(s)
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             "OUTCOMES YOU CAN MEASURE IN 90 DAYS",
             size=11, bold=True, color=BRAND)
    add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.8),
             "We don't ship features. We ship outcomes.",
             size=30, bold=True, color=INK)
    rows = [
        ("Metric", "Today", "Day 90 with Niyamone"),
        ("Report turnaround time", "4.2 hrs avg", "≤ 2.5 hrs avg"),
        ("Manual reconciliation", "Daily, 2 staff", "Zero"),
        ("Untracked discount value", "~ ₹ 35 K / month", "₹ 0 (every discount approved)"),
        ("Branch reconciliation", "End-of-week", "Real-time"),
        ("Patient WhatsApp delivery", "Manual upload", "1-click, with signed PDF"),
    ]
    base_y = Inches(2.4)
    col_x = [Inches(0.6), Inches(5.4), Inches(8.0)]
    col_w = [Inches(4.8), Inches(2.6), Inches(5.0)]
    for i, row in enumerate(rows):
        rh = Inches(0.55)
        fill = BRAND if i == 0 else (WHITE if i % 2 else SURFACE_SUB)
        for j in range(3):
            add_rect(s, col_x[j], base_y + i * rh, col_w[j], rh,
                     fill=fill, line=INK_FAINT, line_w=0.4)
            add_text(s, col_x[j] + Inches(0.18), base_y + i * rh + Inches(0.12),
                     col_w[j] - Inches(0.2), rh,
                     row[j],
                     size=14 if i else 13,
                     bold=bool((i == 0) or (j > 0 and i > 0)),
                     color=WHITE if i == 0 else INK)
    page_footer(s, idx, total)
    return s


def make_roi(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SURFACE)
    brand_top_bar(s)
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             "THE MATH", size=11, bold=True, color=BRAND)
    add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.8),
             "Payback in under 5 months. Then it's pure margin.",
             size=28, bold=True, color=INK)
    cards = [
        ("STATUS QUO / YEAR",  "₹ 8.6 L",
         "2 reconcilers + paper reports + leaked discounts + 3 disjoint tools.", BRAND_DEEP),
        ("NIYAMONE / YEAR",    "₹ 3.2 L",
         "All-in licence, support, 2 training cycles, unlimited branches.", BRAND),
        ("YEAR-1 SAVING",      "₹ 5.4 L",
         "Plus reclaimed staff hours redeployed to revenue work.", GOOD_FG),
    ]
    for i, (k, big, copy, c) in enumerate(cards):
        x = Inches(0.6 + i * 4.25)
        add_rect(s, x, Inches(2.6), Inches(4.05), Inches(3.4),
                 fill=WHITE, line=INK_FAINT, line_w=0.5)
        add_rect(s, x, Inches(2.6), Inches(4.05), Inches(0.6), fill=c)
        add_text(s, x + Inches(0.25), Inches(2.7), Inches(4), Inches(0.4),
                 k, size=11, bold=True, color=WHITE)
        add_text(s, x + Inches(0.25), Inches(3.4), Inches(4), Inches(1.0),
                 big, size=44, bold=True, color=c)
        add_text(s, x + Inches(0.25), Inches(4.6), Inches(3.6), Inches(1.4),
                 copy, size=12, color=INK_MUTED)
    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
             "Figures based on a 5-branch lab benchmark. Actual numbers will be co-modelled in week 1.",
             size=10, color=INK_FAINT)
    page_footer(s, idx, total)
    return s


def make_section_divider(prs, idx, total, section_num: str, title: str, subtitle: str):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=BRAND_DEEP)
    add_rect(s, 0, Inches(6.7), SLIDE_W, Inches(0.08), fill=ACCENT)
    add_text(s, Inches(0.8), Inches(1.4), Inches(6), Inches(2),
             section_num, size=180, bold=True, color=BRAND)
    add_text(s, Inches(5.5), Inches(2.6), Inches(7.5), Inches(0.5),
             "SECTION", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(5.5), Inches(3.1), Inches(7.5), Inches(1.4),
             title, size=44, bold=True, color=WHITE)
    add_text(s, Inches(5.5), Inches(4.7), Inches(7.5), Inches(1.5),
             subtitle, size=18, color=INK_FAINT)
    page_footer(s, idx, total)
    return s


def make_walkthrough(prs, idx, total, *,
                     section_label: str,
                     headline: str,
                     bullets: list[str],
                     screenshot: str,
                     caption: str,
                     saves: str):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SURFACE)
    brand_top_bar(s)
    add_text(s, Inches(0.6), Inches(0.6), Inches(8), Inches(0.4),
             section_label, size=11, bold=True, color=BRAND)
    add_text(s, Inches(0.6), Inches(1.0), Inches(7.0), Inches(1.3),
             headline, size=28, bold=True, color=INK)
    add_bullets(s, Inches(0.6), Inches(2.6), Inches(6.6), Inches(3.4),
                bullets, size=15, color=INK_SOFT, gap=10)

    # "What it saves you" callout
    add_rect(s, Inches(0.6), Inches(5.8), Inches(6.6), Inches(1.0),
             fill=WHITE, line=GOOD_FG, line_w=1.5)
    add_pill(s, Inches(0.8), Inches(5.95), "WHAT IT SAVES YOU",
             fill=GOOD_FG, w=Inches(1.95))
    add_text(s, Inches(0.8), Inches(6.4), Inches(6.4), Inches(0.4),
             saves, size=13, bold=True, color=INK)

    # Screenshot 40% right
    add_screenshot(s, screenshot, Inches(7.6), Inches(1.0),
                   Inches(5.2), Inches(4.8), caption)
    page_footer(s, idx, total)
    return s


def make_trust(prs, idx, total, *, section_label, headline, bullets, screenshot, caption):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SURFACE)
    brand_top_bar(s)
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             section_label, size=11, bold=True, color=BRAND)
    add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(1.0),
             headline, size=30, bold=True, color=INK)
    add_bullets(s, Inches(0.6), Inches(2.8), Inches(6.6), Inches(4),
                bullets, size=15, color=INK_SOFT, gap=12)
    add_screenshot(s, screenshot, Inches(7.6), Inches(2.6),
                   Inches(5.2), Inches(3.6), caption)
    page_footer(s, idx, total)
    return s


def make_rollout(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SURFACE)
    brand_top_bar(s)
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             "HOW WE ROLL THIS OUT", size=11, bold=True, color=BRAND)
    add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.8),
             "Three phases. Fixed milestones. You sign off each one.",
             size=28, bold=True, color=INK)
    phases = [
        ("PHASE 1", "Weeks 1–3",
         "Tenant set-up, master data, lab + billing live at branch HQ.",
         "Go-live: first invoice + first signed report from Niyamone."),
        ("PHASE 2", "Weeks 4–6",
         "Rest of the branches onboarded, pharmacy + home collection on.",
         "Go-live: all branches live, WhatsApp delivery on."),
        ("PHASE 3", "Weeks 7–9",
         "Doctor payouts, smart-inbox approvals, dashboards, training.",
         "Go-live: monthly close run end-to-end inside Niyamone."),
    ]
    for i, (k, w, scope, exit_) in enumerate(phases):
        x = Inches(0.6 + i * 4.25)
        add_rect(s, x, Inches(2.6), Inches(4.05), Inches(4.1),
                 fill=WHITE, line=INK_FAINT, line_w=0.5)
        add_rect(s, x, Inches(2.6), Inches(4.05), Inches(0.7), fill=BRAND)
        add_text(s, x + Inches(0.25), Inches(2.7), Inches(4), Inches(0.4),
                 k, size=12, bold=True, color=WHITE)
        add_text(s, x + Inches(0.25), Inches(3.0), Inches(4), Inches(0.35),
                 w, size=11, color=ACCENT, bold=True)
        add_text(s, x + Inches(0.25), Inches(3.6), Inches(3.6), Inches(1.4),
                 scope, size=14, color=INK_SOFT)
        add_pill(s, x + Inches(0.25), Inches(5.2), "EXIT CRITERIA",
                 fill=GOLD, fg=INK, w=Inches(1.6))
        add_text(s, x + Inches(0.25), Inches(5.65), Inches(3.6), Inches(1.0),
                 exit_, size=12, bold=True, color=INK)
    page_footer(s, idx, total)
    return s


def make_close(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=BRAND_DEEP)
    add_rect(s, 0, Inches(0.0), SLIDE_W, Inches(0.12), fill=ACCENT)
    add_text(s, Inches(0.8), Inches(0.8), Inches(12), Inches(0.5),
             "TODAY WE AGREE", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), Inches(1.4), Inches(12), Inches(1.6),
             "Three signatures. Nine weeks to live.",
             size=44, bold=True, color=WHITE)

    boxes = [
        ("1. Scope locked",  "Phases, branches, integrations.",
         "Sign-off by:  ___________________"),
        ("2. Commercials",   "Licence + implementation + support.",
         "Sign-off by:  ___________________"),
        ("3. Kick-off date", "Week 1 starts Monday.",
         "Date:          ___________________"),
    ]
    for i, (k, sub, sig) in enumerate(boxes):
        x = Inches(0.8 + i * 4.2)
        add_rect(s, x, Inches(3.4), Inches(4.0), Inches(3.0),
                 fill=WHITE, line=ACCENT, line_w=1.5)
        add_text(s, x + Inches(0.25), Inches(3.55), Inches(3.7), Inches(0.5),
                 k, size=18, bold=True, color=BRAND)
        add_text(s, x + Inches(0.25), Inches(4.1), Inches(3.7), Inches(0.9),
                 sub, size=13, color=INK_MUTED)
        add_text(s, x + Inches(0.25), Inches(5.4), Inches(3.7), Inches(0.9),
                 sig, size=12, color=INK)
    add_text(s, Inches(0.8), Inches(6.8), Inches(12), Inches(0.4),
             "Venkki M. K.  ·  Niyamone Labs  ·  venkki.mk@niyamone.com  ·  +91-XXXXXXXXXX",
             size=11, color=INK_FAINT)
    return s


# ── Deck definition ───────────────────────────────────────────────────
WALKTHROUGH = [
    # (section_label, headline, bullets, screenshot, caption, saves)
    ("ACT III  ·  01 / 13  ·  LOGIN", "Role-based login — the same app, scoped to each user.",
     ["JWT carries staff_id + branch_id + role; nothing is trusted from the URL.",
      "Super-admin sees all branches; branch-admin sees one; reception sees their own queue.",
      "Inactive staff are blocked at the auth layer, not at the screen."],
     "01-login.png",
     "Reception logs in once, lands on her branch — never sees the other four.",
     "Zero accidental cross-branch edits. Every action is name-stamped."),

    ("ACT III  ·  02 / 13  ·  DASHBOARD", "One screen — every KPI, every branch.",
     ["Revenue, TAT, pending verifications, category share — live.",
      "Branch slicer at the top reshapes every tile in <200 ms.",
      "Export to PDF matches the on-screen design — pixel-for-pixel."],
     "02-dashboard.png",
     "Owner sees the day at 8 AM; doesn't need to call anyone.",
     "End-of-day reconciliation goes from 45 min to instant."),

    ("ACT III  ·  03 / 13  ·  PATIENTS", "Register a patient in under 20 seconds.",
     ["Duplicate guard on mobile + name + DOB combo — no UHID drift.",
      "Address book reused across home collection and report delivery.",
      "Gender defaults to Male so reception can tab through and hit Save."],
     "03-patient-register.png",
     "New patient, UHID, mobile, address — one screen, one save.",
     "A 4-minute paper form becomes a 20-second tap-through."),

    ("ACT III  ·  04 / 13  ·  APPOINTMENTS", "Home collection that doesn't lose ₹250.",
     ["Pickup surcharge auto-added to the invoice — editable + discountable.",
      "Address captured once, reused on report + WhatsApp share.",
      "Scheduled-at time shows on the invoice's visit-details block."],
     "04-home-collection.png",
     "Surcharge is a line item, not a sticky note on the cash register.",
     "Eliminates the most common cause of cash leakage in home collection."),

    ("ACT III  ·  05 / 13  ·  BILLING", "Invoice that knows where every line came from.",
     ["Lab, pharmacy, doctor, IPD, manual — each line tagged with its provenance.",
      "GST + discount tiers + payment + balance — all on one screen.",
      "WhatsApp share with public link + auto-print on the patient's phone."],
     "05-billing-invoice.png",
     "Every paisa traces back to a clinical action. No mystery charges.",
     "Audit-ready billing the first day you go live."),

    ("ACT III  ·  06 / 13  ·  SMART INBOX", "No discount goes through unapproved.",
     ["Tiered: auto / branch-admin / super-admin — based on the percentage.",
      "Submitted requests appear in the approver's Smart Inbox in real time.",
      "Approval, rejection, and apply-error are all audit-trailed."],
     "06-smart-inbox.png",
     "Discount approval becomes a 30-second chat, not a phone call.",
     "Plugs the single biggest revenue leak at the cash counter."),

    ("ACT III  ·  07 / 13  ·  LAB WORKFLOW", "Accession → sample → result → verify → report.",
     ["One queue per stage. Items can't skip stages.",
      "Critical alerts surface immediately to the doctor on duty.",
      "Verifier sign-off captured via uploaded digital signature."],
     "07-lab-workflow.png",
     "Lab tech, pathologist, doctor — same workflow, three perspectives.",
     "TAT visibility turns from a guess into a real-time number."),

    ("ACT III  ·  08 / 13  ·  LAB REPORT PDF", "The Sree Diagnostics letterhead — pixel-perfect.",
     ["Header, footer, accreditations, seals, watermark — all configurable.",
      "QR code on the footer points to a verifiable public URL.",
      "Filename is PatientName_DD-MMM-YYYY_UHID — searchable in WhatsApp."],
     "08-lab-report-pdf.png",
     "The exact PDF you sent us — now generated in two clicks.",
     "Brand consistency across every report, every branch, every device."),

    ("ACT III  ·  09 / 13  ·  PHARMACY", "Dispense, indent, GRN — all in flow.",
     ["Stock-aware dispensing — won't sell what isn't on the shelf.",
      "Expiry guard at dispense time — flagged in red.",
      "Sales auto-flow to the invoice with HSN + GST."],
     "09-pharmacy.png",
     "Pharmacy is no longer a parallel universe; it's part of the bill.",
     "Stops the daily 'pharmacy vs billing' tally meeting."),

    ("ACT III  ·  10 / 13  ·  IPD / WARDS", "Bed map, doctor visits, consolidated bill.",
     ["Bed assignments time-tracked for accurate per-day billing.",
      "Doctor visits roll up into the same invoice the lab uses.",
      "Discharge produces one PDF: bill + summary + reports."],
     "10-ipd.png",
     "One discharge — one PDF — one signature.",
     "Discharge bottleneck disappears."),

    ("ACT III  ·  11 / 13  ·  DOCTOR PAYOUTS", "Indian-standard payslip, in two clicks.",
     ["Earnings + deductions side-by-side; TDS Section 194J calculated.",
      "Amount-in-words in Lakhs / Crores — auditor-friendly.",
      "Lab header, branch address, GST — same letterhead as the report."],
     "11-doctor-payslip.png",
     "Payslip for a consultant takes 6 seconds, not 6 minutes.",
     "Month-end payouts cease to be a spreadsheet rodeo."),

    ("ACT III  ·  12 / 13  ·  REPORTS & ANALYTICS", "On-screen and on-paper — identical.",
     ["Category share, revenue trend, branch drill — all interactive.",
      "PDF export uses the same DOM, not a separate template.",
      "Date range, branch, category — three controls, no SQL."],
     "12-reports.png",
     "What you see is what you print — full stop.",
     "No more 'why does the PDF look different from the screen?'"),

    ("ACT III  ·  13 / 13  ·  SETTINGS", "Your brand, your seals, your instructions.",
     ["Logo, accreditations, watermark, footer seals — all uploadable.",
      "Instructions per test / per branch / global — cascades intelligently.",
      "Per-staff digital signature stored once, used everywhere."],
     "13-settings.png",
     "The branding doesn't need a developer; reception can change it.",
     "Re-branding moments (new NABL cert, new address) take 5 minutes."),
]

TRUST = [
    ("ARCHITECTURE", "Your data lives in your tenant.",
     ["Angular 21 frontend, Supabase backend, RLS per table.",
      "All writes go through SECURITY DEFINER RPCs — no rogue UPDATEs possible.",
      "Database snapshots are yours, exportable on demand, in standard Postgres dump.",
      "No vendor lock-in: walk away with a .sql file."],
     "14-architecture.png",
     "Schema diagram — every table owned by the client's Supabase project."),
    ("SECURITY & COMPLIANCE", "Row-level security, not screen-level pretending.",
     ["Postgres RLS gates every read and write by branch + role.",
      "JWT custom claims (app_metadata) carry the role — tamper-proof.",
      "Audit trail on every exception (discounts, voids, refunds).",
      "Storage buckets for PDFs are private with signed URLs."],
     "15-security.png",
     "RLS policies block cross-branch reads at the database, not the UI."),
    ("MULTI-BRANCH", "Switch branches in one click — no relogin.",
     ["BranchStore is reactive; every screen reshapes on switch.",
      "Realtime channels keep dashboards live across all branches.",
      "Reports and exports stamp the active branch in the footer."],
     "16-multi-branch.png",
     "Top-bar branch switcher — propagates to every page instantly."),
]


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Total = 25 slides:
    # 1 cover + 3 act-I + 3 act-II + 1 act-III divider + 13 walkthrough
    # + 1 act-IV divider + 3 trust + 1 rollout + 1 close
    total = 27

    idx = 1
    make_cover(prs);                                                          idx_ = 1
    make_value_prop(prs, 2, total)
    make_problem(prs,    3, total)
    make_promise(prs,    4, total)
    make_outcomes(prs,   5, total)
    make_roi(prs,        6, total)

    make_section_divider(prs, 7, total, "01",
                         "The walkthrough.",
                         "Thirteen screens. Thirteen reasons this is already running for you.")
    n = 8
    for w in WALKTHROUGH:
        make_walkthrough(prs, n, total,
                         section_label=w[0], headline=w[1], bullets=w[2],
                         screenshot=w[3], caption=w[4], saves=w[5])
        n += 1

    make_section_divider(prs, n, total, "02",
                         "Built for trust.",
                         "Architecture, security, and the multi-branch reality.")
    n += 1
    for t in TRUST:
        make_trust(prs, n, total,
                   section_label=t[0], headline=t[1], bullets=t[2],
                   screenshot=t[3], caption=t[4])
        n += 1

    make_rollout(prs, n, total); n += 1
    make_close(prs,   n, total); n += 1

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"OK  Deck written: {OUT}")
    print(f"  Slides: {n - 1}")
    print(f"  Drop screenshots into: {SHOTS}")


if __name__ == "__main__":
    build()
